"""
Compliance Assessment PPTX generator — Suricatoos Corporativo Navy theme.

Slides:
  1  Cover          – framework, assessment name, date, confidential badge
  2  Sumário        – 5 KPI cards (compliant / partial / non / N/A / unassessed)
  3  Conformidade   – donut chart + compliance score
  4  Por Domínio    – spider/radar chart
  5  Progresso      – completion bar chart per category
  6  Controles      – applied controls status bar
  7  P1 Controls    – priority-1 controls table
  8  Não Conformes  – non-compliant/partial requirements table (up to 20 rows)
  9  Closing        – Obrigado / Merci / Thank you
"""
import io
from math import ceil

from django.db.models import Count
from django.utils.timezone import now

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from library.helpers import get_referential_translation

from .generators import (
    plot_donut,
    plot_spider_chart,
    plot_completion_bar,
    plot_horizontal_bar,
)
from .models import (
    ComplianceAssessment,
    RequirementAssessment,
    AppliedControl,
)


# ─── Corporativo Navy palette ─────────────────────────────────────────────────
NAVY    = RGBColor(0x0f, 0x17, 0x2a)   # slate-950  cover bg
SLATE8  = RGBColor(0x1e, 0x29, 0x3b)   # slate-800  header bars
SLATE7  = RGBColor(0x33, 0x41, 0x55)   # slate-700  body text
SLATE5  = RGBColor(0x64, 0x74, 0x8b)   # slate-500  muted text
SLATE4  = RGBColor(0x94, 0xa3, 0xb8)   # slate-400  secondary labels
SLATE1  = RGBColor(0xf1, 0xf5, 0xf9)   # slate-100  alt row bg
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x3b, 0x82, 0xf6)   # blue-500
EMERALD = RGBColor(0x10, 0xb9, 0x81)   # emerald-500  compliant
AMBER   = RGBColor(0xf5, 0x9e, 0x0b)   # amber-500    partially compliant
ROSE    = RGBColor(0xef, 0x44, 0x44)   # red-500      non-compliant
GRAY    = RGBColor(0x64, 0x74, 0x8b)   # slate-500    N/A / not assessed
ROSE_BG = RGBColor(0xff, 0xf1, 0xf0)   # light rose   non-compliant row
AMB_BG  = RGBColor(0xff, 0xfd, 0xed)   # light amber  partial row

CHART_BLUE = ["#3b82f6"]
CHART_AC   = ["#e2e8f0", "#10b981", "#ef4444", "#334155", "#f59e0b", "#bfdbfe"]

W = Inches(13.33)
H = Inches(7.5)
M = Inches(0.5)


# ─── i18n ─────────────────────────────────────────────────────────────────────
_L = {
    "pt": {
        "compliant": "Conforme",
        "partially_compliant": "Parcialmente conforme",
        "non_compliant": "Não conforme",
        "not_applicable": "Não aplicável",
        "not_assessed": "Não avaliado",
        "exec_summary": "Resumo Executivo",
        "compliance_overview": "Visão Geral de Conformidade",
        "by_domain": "Conformidade por Domínio",
        "by_category": "Progresso por Categoria",
        "applied_controls": "Controles Aplicados",
        "p1_controls": "Controles de Alta Prioridade (P1)",
        "non_compliant_reqs": "Requisitos Não Conformes",
        "closing": "Obrigado",
        "total_reqs": "Total de requisitos",
        "compliance_score": "Índice de conformidade",
        "confidential": "CONFIDENCIAL",
        "prepared_by": "Elaborado por",
        "ref": "Ref", "name": "Nome", "status": "Status",
        "category": "Categoria", "coverage": "Cobertura", "controls_col": "Controles",
        "to_do": "A fazer", "in_progress": "Em andamento", "active": "Ativo",
        "done": "Concluído", "on_hold": "Em espera", "deprecated": "Obsoleto",
        "in_review": "Em revisão", "policy": "Política", "process": "Processo",
        "technical": "Técnico", "physical": "Físico", "procedure": "Procedimento",
    },
    "fr": {
        "compliant": "Conformes",
        "partially_compliant": "Partiellement conformes",
        "non_compliant": "Non conformes",
        "not_applicable": "Non applicables",
        "not_assessed": "Non évalués",
        "exec_summary": "Résumé Exécutif",
        "compliance_overview": "Vue d'ensemble",
        "by_domain": "Conformité par Domaine",
        "by_category": "Complétion par Catégorie",
        "applied_controls": "Contrôles Appliqués",
        "p1_controls": "Contrôles Priorité 1",
        "non_compliant_reqs": "Exigences Non Conformes",
        "closing": "Merci",
        "total_reqs": "Total exigences",
        "compliance_score": "Score de conformité",
        "confidential": "CONFIDENTIEL",
        "prepared_by": "Préparé par",
        "ref": "Réf", "name": "Nom", "status": "Statut",
        "category": "Catégorie", "coverage": "Couverture", "controls_col": "Contrôles",
        "to_do": "À faire", "in_progress": "En cours", "active": "Actif",
        "done": "Terminé", "on_hold": "En attente", "deprecated": "Déprécié",
        "in_review": "En revue", "policy": "Politique", "process": "Processus",
        "technical": "Technique", "physical": "Physique", "procedure": "Procédure",
    },
    "en": {
        "compliant": "Compliant",
        "partially_compliant": "Partially compliant",
        "non_compliant": "Non-compliant",
        "not_applicable": "Not applicable",
        "not_assessed": "Not assessed",
        "exec_summary": "Executive Summary",
        "compliance_overview": "Compliance Overview",
        "by_domain": "Compliance by Domain",
        "by_category": "Completion by Category",
        "applied_controls": "Applied Controls",
        "p1_controls": "Priority 1 Controls",
        "non_compliant_reqs": "Non-Compliant Requirements",
        "closing": "Thank you",
        "total_reqs": "Total requirements",
        "compliance_score": "Compliance Score",
        "confidential": "CONFIDENTIAL",
        "prepared_by": "Prepared by",
        "ref": "Ref", "name": "Name", "status": "Status",
        "category": "Category", "coverage": "Coverage", "controls_col": "Controls",
        "to_do": "To do", "in_progress": "In progress", "active": "Active",
        "done": "Done", "on_hold": "On hold", "deprecated": "Deprecated",
        "in_review": "In review", "policy": "Policy", "process": "Process",
        "technical": "Technical", "physical": "Physical", "procedure": "Procedure",
    },
}


def _t(lang: str, key: str) -> str:
    return _L.get(lang, _L["en"]).get(key, _L["en"].get(key, key))


def _tv(lang: str, val) -> str:
    if val is None or val == "--":
        return "-"
    s = str(val)
    return _t(lang, s) if s in _L.get(lang, _L["en"]) else s


# ─── python-pptx helpers ──────────────────────────────────────────────────────

def _bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height, text, size=12,
         bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def _rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor = None):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def _img(slide, buf: io.BytesIO, left, top, width=None, height=None):
    buf.seek(0)
    if width and height:
        slide.shapes.add_picture(buf, left, top, width, height)
    elif width:
        slide.shapes.add_picture(buf, left, top, width=width)
    elif height:
        slide.shapes.add_picture(buf, left, top, height=height)
    else:
        slide.shapes.add_picture(buf, left, top)


def _header(slide, title: str):
    _rect(slide, 0, 0, W, Inches(0.8), SLATE8)
    _box(slide, M, Inches(0.1), W - Inches(3.2), Inches(0.6), title, 18, bold=True)
    _box(slide, W - Inches(2.9), Inches(0.13), Inches(2.7), Inches(0.5),
         "SURICATOOS vCISO", 9, color=SLATE4, align=PP_ALIGN.RIGHT)


# ─── slide builders ───────────────────────────────────────────────────────────

def _cover(prs, audit, contributors: str, lang: str):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, NAVY)
    _rect(sl, 0, 0, Inches(0.12), H, BLUE)
    _box(sl, W - Inches(3.8), Inches(0.32), Inches(3.5), Inches(0.5),
         "SURICATOOS vCISO", 12, color=BLUE, align=PP_ALIGN.RIGHT)
    fw = audit.framework.name if audit.framework else ""
    _box(sl, Inches(0.6), Inches(1.55), Inches(11.5), Inches(1.1),
         fw, 38, bold=True, color=WHITE)
    _box(sl, Inches(0.6), Inches(2.8), Inches(11.5), Inches(0.65),
         audit.name, 22, color=SLATE4)
    if audit.perimeter:
        _box(sl, Inches(0.6), Inches(3.55), Inches(10), Inches(0.45),
             str(audit.perimeter), 14, color=SLATE5)
    if contributors:
        _box(sl, Inches(0.6), Inches(4.25), Inches(9), Inches(0.45),
             f"{_t(lang, 'prepared_by')}: {contributors}", 11, color=SLATE5)
    _box(sl, Inches(0.6), Inches(5.1), Inches(5), Inches(0.38),
         now().strftime("%d/%m/%Y"), 12, color=SLATE4)
    _rect(sl, Inches(9.8), Inches(5.05), Inches(3.0), Inches(0.5), BLUE)
    _box(sl, Inches(9.8), Inches(5.05), Inches(3.0), Inches(0.5),
         _t(lang, "confidential"), 13, bold=True, align=PP_ALIGN.CENTER)


def _exec_summary(prs, agg: dict, lang: str):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, WHITE)
    _header(sl, _t(lang, "exec_summary"))
    total = agg.get("total", 1) or 1
    cards = [
        (_t(lang, "compliant"),           agg.get("compliant", 0),           EMERALD),
        (_t(lang, "partially_compliant"), agg.get("partially_compliant", 0), AMBER),
        (_t(lang, "non_compliant"),       agg.get("non_compliant", 0),       ROSE),
        (_t(lang, "not_applicable"),      agg.get("not_applicable", 0),      GRAY),
        (_t(lang, "not_assessed"),        agg.get("not_assessed", 0),        BLUE),
    ]
    cw, ch, gap = Inches(2.3), Inches(2.2), Inches(0.17)
    sx = (W - (5 * cw + 4 * gap)) / 2
    sy = Inches(1.45)
    for i, (label, count, color) in enumerate(cards):
        x = sx + i * (cw + gap)
        _rect(sl, x, sy, cw, ch, color)
        _box(sl, x, sy + Inches(0.15), cw, Inches(0.6),
             str(count), 36, bold=True, align=PP_ALIGN.CENTER)
        pct = ceil(count / total * 100)
        _box(sl, x, sy + Inches(0.78), cw, Inches(0.35),
             f"{pct}%", 16, align=PP_ALIGN.CENTER)
        _box(sl, x, sy + Inches(1.18), cw, Inches(0.85),
             label, 10, align=PP_ALIGN.CENTER, wrap=True)
    _box(sl, M, Inches(4.25), W - 2*M, Inches(0.38),
         f"{_t(lang, 'total_reqs')}: {agg.get('total', 0)}",
         12, color=SLATE7, align=PP_ALIGN.CENTER)


def _donut_slide(prs, donut_buf: io.BytesIO, agg: dict, lang: str):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, WHITE)
    _header(sl, _t(lang, "compliance_overview"))
    _img(sl, donut_buf, M, Inches(0.9), width=Inches(7.8))
    total = agg.get("total", 1) or 1
    ok = agg.get("compliant", 0) + agg.get("not_applicable", 0)
    pct = ceil(ok / total * 100)
    _box(sl, Inches(8.4), Inches(1.6), Inches(4.5), Inches(0.6),
         _t(lang, "compliance_score"), 15, color=SLATE7)
    _box(sl, Inches(8.4), Inches(2.3), Inches(4.5), Inches(1.6),
         f"{pct}%", 64, bold=True, color=BLUE)


def _spider_slide(prs, spider_buf: io.BytesIO, lang: str):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, WHITE)
    _header(sl, _t(lang, "by_domain"))
    _img(sl, spider_buf, M, Inches(0.9), height=Inches(6.3))


def _completion_slide(prs, completion_buf: io.BytesIO, lang: str):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, WHITE)
    _header(sl, _t(lang, "by_category"))
    _img(sl, completion_buf, M, Inches(0.9), width=W - 2*M)


def _controls_slide(prs, ac_buf: io.BytesIO, ac_count: int, lang: str):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, WHITE)
    _header(sl, _t(lang, "applied_controls"))
    _img(sl, ac_buf, M, Inches(0.9), width=Inches(10.5))
    _box(sl, Inches(11.0), Inches(1.5), Inches(2.0), Inches(0.5),
         f"Total: {ac_count}", 14, color=SLATE7)


def _p1_slide(prs, p1: list, lang: str):
    if not p1:
        return
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, WHITE)
    _header(sl, _t(lang, "p1_controls"))
    cols = [
        (_t(lang, "name"),     Inches(4.5)),
        (_t(lang, "category"), Inches(2.0)),
        (_t(lang, "status"),   Inches(2.2)),
        (_t(lang, "coverage"), Inches(1.3)),
    ]
    rh = Inches(0.38)
    y = Inches(0.87)
    x = M
    for hdr, w in cols:
        _rect(sl, x, y, w, rh, SLATE8)
        _box(sl, x + Inches(0.05), y + Inches(0.04), w - Inches(0.1), rh, hdr, 10, bold=True)
        x += w
    for i, ctrl in enumerate(p1[:14]):
        y += rh
        x = M
        bg = SLATE1 if i % 2 == 0 else WHITE
        row = [
            (ctrl.get("name") or "-")[:60],
            _tv(lang, ctrl.get("category")),
            _tv(lang, ctrl.get("status")),
            str(ctrl.get("coverage", 0)),
        ]
        for val, (_, w) in zip(row, cols):
            _rect(sl, x, y, w, rh, bg)
            _box(sl, x + Inches(0.05), y + Inches(0.03), w - Inches(0.1), rh, val, 9, color=SLATE7)
            x += w


def _non_compliant_slide(prs, ra_list: list, lang: str):
    bad = [
        ra for ra in ra_list
        if ra.get("result") in ("non_compliant", "partially_compliant")
    ][:20]
    if not bad:
        return
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, WHITE)
    _header(sl, _t(lang, "non_compliant_reqs"))
    cols = [
        (_t(lang, "ref"),          Inches(0.95)),
        (_t(lang, "name"),         Inches(4.5)),
        (_t(lang, "status"),       Inches(2.0)),
        (_t(lang, "controls_col"), Inches(5.38)),
    ]
    rh = Inches(0.33)
    y = Inches(0.87)
    x = M
    for hdr, w in cols:
        _rect(sl, x, y, w, rh, SLATE8)
        _box(sl, x + Inches(0.04), y + Inches(0.03), w - Inches(0.08), rh, hdr, 10, bold=True)
        x += w
    for ra in bad:
        y += rh
        x = M
        result = ra.get("result", "")
        row_bg = ROSE_BG if result == "non_compliant" else AMB_BG
        row = [
            ra.get("ref_id", "-"),
            (ra.get("name") or "-")[:60],
            _tv(lang, result),
            (ra.get("applied_controls") or "-")[:80],
        ]
        for val, (_, w) in zip(row, cols):
            _rect(sl, x, y, w, rh, row_bg)
            _box(sl, x + Inches(0.04), y + Inches(0.02), w - Inches(0.08), rh, val, 8, color=SLATE7)
            x += w


def _closing(prs, lang: str):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, NAVY)
    _rect(sl, 0, 0, Inches(0.12), H, BLUE)
    _box(sl, W - Inches(3.8), Inches(0.32), Inches(3.5), Inches(0.5),
         "SURICATOOS vCISO", 12, color=BLUE, align=PP_ALIGN.RIGHT)
    _box(sl, Inches(1.0), Inches(2.4), W - Inches(2.0), Inches(1.5),
         _t(lang, "closing"), 52, bold=True, align=PP_ALIGN.CENTER)
    _box(sl, Inches(1.0), Inches(4.1), W - Inches(2.0), Inches(0.45),
         "vciso.suricatoos.com", 13, color=SLATE4, align=PP_ALIGN.CENTER)


# ─── data helpers ─────────────────────────────────────────────────────────────

def _count_results(node_data: dict) -> dict:
    counts: dict = {}
    if node_data.get("assessable"):
        r = node_data.get("result", "unknown")
        counts[r] = 1
    for child in node_data.get("children", {}).values():
        for k, v in _count_results(child).items():
            counts[k] = counts.get(k, 0) + v
    return counts


# ─── main entry ───────────────────────────────────────────────────────────────

def build_compliance_pptx(audit_id: str, tree: dict, lang: str = "pt") -> io.BytesIO:
    """Build a compliance assessment PPTX presentation. Returns a BytesIO buffer."""
    audit = ComplianceAssessment.objects.get(id=audit_id)

    # aggregate compliance counts
    result_counts = {nd["urn"]: _count_results(nd) for nd in tree.values()}
    agg = {
        "compliant": 0, "non_compliant": 0,
        "not_applicable": 0, "not_assessed": 0, "partially_compliant": 0,
    }
    for counts in result_counts.values():
        for k in agg:
            agg[k] += counts.get(k, 0)
    agg["total"] = sum(agg.values())

    # spider / completion data (% compliant+N/A per top-level domain)
    spider_data = []
    for nd in tree.values():
        counts = result_counts.get(nd["urn"], {})
        n = sum(counts.values())
        ok = counts.get("compliant", 0) + counts.get("not_applicable", 0)
        pct = ceil(ok / n * 100) if n > 0 else 0
        spider_data.append({"category": nd["node_content"].split(":")[0], "value": pct})

    donut_data = [
        {"category": _t(lang, "compliant"),           "value": agg["compliant"]},
        {"category": _t(lang, "partially_compliant"), "value": agg["partially_compliant"]},
        {"category": _t(lang, "non_compliant"),       "value": agg["non_compliant"]},
        {"category": _t(lang, "not_applicable"),      "value": agg["not_applicable"]},
        {"category": _t(lang, "not_assessed"),        "value": agg["not_assessed"]},
    ]

    # applied controls
    req_qs = audit.get_requirement_assessments(include_non_assessable=True)
    controls = AppliedControl.objects.filter(
        requirement_assessments__in=req_qs
    ).distinct()
    ac_count = controls.count()
    ac_status = controls.values("status").annotate(n=Count("id"))
    ac_data = [
        {"category": _tv(lang, s["status"]), "value": s["n"]}
        for s in ac_status
    ]
    p1 = []
    for ac in controls.filter(priority=1):
        cov = RequirementAssessment.objects.filter(
            compliance_assessment=audit, applied_controls=ac.id
        ).count()
        p1.append({"name": ac.name, "status": ac.status,
                   "category": ac.category, "coverage": cov})

    # requirement list
    ra_list = []
    for ra in req_qs:
        if not ra.requirement.assessable:
            continue
        ra_list.append({
            "ref_id": ra.requirement.ref_id or "-",
            "name": get_referential_translation(ra.requirement, "name", lang) or "-",
            "result": ra.result,
            "applied_controls": ", ".join(ac.name for ac in ra.applied_controls.all()) or "-",
        })

    # contributors
    authors = ", ".join(dict.fromkeys(
        e for a in audit.authors.all() for e in a.get_emails()
    ))
    reviewers = ", ".join(dict.fromkeys(
        e for r in audit.reviewers.all() for e in r.get_emails()
    ))
    contributors = "\n".join(filter(None, [authors, reviewers]))

    # generate charts
    donut_buf      = plot_donut(donut_data)
    spider_buf     = plot_spider_chart(spider_data, colors=CHART_BLUE)
    completion_buf = plot_completion_bar(spider_data, colors=CHART_AC)
    ac_buf         = plot_horizontal_bar(ac_data, colors=CHART_AC) if ac_data else None

    # build presentation
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _cover(prs, audit, contributors, lang)
    _exec_summary(prs, agg, lang)
    _donut_slide(prs, donut_buf, agg, lang)
    _spider_slide(prs, spider_buf, lang)
    _completion_slide(prs, completion_buf, lang)
    if ac_buf:
        _controls_slide(prs, ac_buf, ac_count, lang)
    if p1:
        _p1_slide(prs, p1, lang)
    _non_compliant_slide(prs, ra_list, lang)
    _closing(prs, lang)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out
